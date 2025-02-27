# https://numpy.org/
import numpy as np

# https://pypi.org/project/pandas/
import pandas as pd

# https://docs.python.org/3/library/os.html
import os

# https://docs.python.org/3/library/time.html
import time

# https://streamlit.io/
# https://docs.streamlit.io/develop/api-reference/text
import streamlit as st

# https://www.langchain.com/
from langchain.text_splitter import RecursiveCharacterTextSplitter

# https://python.langchain.com/docs/integrations/text_embedding/openai/
from langchain_openai import OpenAIEmbeddings

# https://faiss.ai/
# https://faiss.ai/cpp_api/struct/structfaiss_1_1IndexFlatL2.html
import faiss

# https://github.com/openai/openai-python
from openai import OpenAI

# https://pypdf2.readthedocs.io/en/3.x/
from PyPDF2 import PdfReader

# https://python-docx.readthedocs.io/en/latest/
import docx

# https://python-pptx.readthedocs.io/en/latest/
import pptx

# https://openpyxl.readthedocs.io/en/stable/
import openpyxl

# https://pypi.org/project/pandoc/
import pypandoc

# https://www.crummy.com/software/BeautifulSoup/bs4/doc/
from bs4 import BeautifulSoup

# https://docs.python.org/3/library/json.html
import json

# https://docs.python.org/3/library/xml.etree.elementtree.html
import xml.etree.ElementTree as ET

# Configuration
# https://docs.streamlit.io/develop/api-reference/configuration/config.toml


def main():
    # Application header
    st.set_page_config(page_title="Query Documents", page_icon="🙋‍♂️", layout="centered", menu_items={'About': 'https://abiro.com'})   
    st.title("📄 Query Documents")
    st.markdown("Upload your documents (PDF, RTF, Word, PowerPoint, Excel, Text, Markdown, HTML, XML or JSON), and ask questions based on their content. The system will only use information from the uploaded files to answer your queries. The provided API Key, set via OPENAI_API_KEY or manually, is only used for queries and not otherwise stored.")

    # OpenAI API key: via the OPENAI_API_KEY environment variable or input
    openai_api_key = os.getenv("OPENAI_API_KEY", "")
    if openai_api_key == "":
        st.header("🔑 OpenAI API Key")
        openai_api_key = st.text_input("Enter your API key (not shown):", type="password")
        st.button("Set")

    if openai_api_key == "":
        st.write("Enter a valid OpenAI API Key")
        return

    # Initialize session state for embeddings and FAISS index
    if 'embeddings' not in st.session_state:
        st.session_state.embeddings = OpenAIEmbeddings(model="text-embedding-3-small", openai_api_key=openai_api_key)

    if 'faiss_index' not in st.session_state:
        init_faiss_index()

    if 'texts' not in st.session_state:
        st.session_state.texts = []

    if 'metadatas' not in st.session_state:
        st.session_state.metadatas = []

    if 'processed_files' not in st.session_state:
        st.session_state.processed_files = set()

    # Sidebar for uploading files
    st.sidebar.header("Upload Files")
    uploaded_files = st.sidebar.file_uploader("Choose files", type=["pdf", "txt", "md", "rtf", "docx", "pptx", "xlsx", "html", "htm", "json", "xml"], accept_multiple_files=True)

    if (len(uploaded_files) < len(st.session_state.processed_files)):
        st.session_state.processed_files = set()        
        init_faiss_index()

    if uploaded_files:
        for uploaded_file in uploaded_files:
            if uploaded_file.name not in st.session_state.processed_files:
                st.sidebar.write(f"📄 {uploaded_file.name}")
                # Extract text based on file type
                text = extract_text_from_file(uploaded_file)
                if text:
                    # Split text into chunks
                    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
                    chunks = text_splitter.split_text(text)
                    add_texts_to_faiss(chunks, st.session_state.embeddings)
                    st.sidebar.success(f"Processed {uploaded_file.name} successfully!")
                    st.session_state.processed_files.add(uploaded_file.name)
                else:
                    st.sidebar.error(f"Failed to extract text from {uploaded_file.name}.")
            else:
                st.sidebar.write(f"📄 {uploaded_file.name} (already processed)")

    # Optional: Show number of indexed chunks
    st.sidebar.write(f"**Total Indexed Chunks:** {st.session_state.faiss_index.ntotal}")

    # Bail if no files uploaded
    if st.session_state.faiss_index.ntotal == 0:
        st.info("Please upload files to start querying.")
        return

    # Query input
    st.header("🙋‍♂️ Question")
    user_query = st.text_input("Type your question here:")
    verbose = st.checkbox("Verbose", help="Check this if you want intermediate information")

    if st.button("Query") and user_query:
        with st.spinner("Querying..."):
            client = OpenAI(api_key=openai_api_key)
            answer = generate_answer(client, user_query, verbose, st.session_state.embeddings, st.session_state.faiss_index, st.session_state.texts)
        
        st.header("👨‍🏫 Answer")
        st.write(answer)

# Initialize FAISS index with the dimensionality from embeddings
def init_faiss_index():
    embedding_test = st.session_state.embeddings.embed_query("test")
    embedding_dim = len(embedding_test)
    st.session_state.faiss_index = faiss.IndexFlatL2(embedding_dim)

# Extract text from a file based on its type
def extract_text_from_file(file) -> str:
    try:
        if file.type == "application/pdf":
            return extract_text_from_pdf(file)
        elif file.type == "text/plain":
            return file.getvalue().decode("utf-8")
        elif file.type == "text/markdown":
            return file.getvalue().decode("utf-8")
        elif file.type == "application/rtf":
            return pypandoc.convert_text(file.getvalue().decode("utf-8"), 'plain', format='rtf')
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return extract_text_from_docx(file)
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return extract_text_from_pptx(file)
        elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return extract_text_from_xlsx(file)
        elif file.type == "text/html":
            return extract_text_from_html(file)
        elif file.type == "application/json":
            return extract_text_from_json(file)
        elif file.type == "application/xml" or file.type == "text/xml":
            return extract_text_from_xml(file)
        elif file.type == "application/octet-stream":
            return file.getvalue().decode("utf-8")       
        else:
            return ""
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""


# Extract text from a PDF file.
def extract_text_from_pdf(file) -> str:
    try:
        pdf = PdfReader(file)
        text = ""
        for page_num, page in enumerate(pdf.pages):
            # TODO Remove header and footer
            page_text = page.extract_text()
            if page_text:
                text += page_text
        return text
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""


# Extract text from a Word document.
def extract_text_from_docx(file) -> str:
    try:
        doc = docx.Document(file)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""


# Extract text from a PowerPoint presentation.
def extract_text_from_pptx(file) -> str:
    try:
        presentation = pptx.Presentation(file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""


# Extract text from an Excel file.
def extract_text_from_xlsx(file) -> str:
    try:
        workbook = openpyxl.load_workbook(file, data_only=True)
        text = ""
        for sheet in workbook.sheetnames:
            worksheet = workbook[sheet]
            for row in worksheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                text += "\t".join(row_text) + "\n"
        return text
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""


# Extract text from an HTML file.
def extract_text_from_html(file) -> str:
    try:
        html_content = file.getvalue().decode("utf-8")
        soup = BeautifulSoup(html_content, "html.parser")
        text = soup.get_text(separator="\n")
        return text
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""


# Extract text from a JSON file.
def extract_text_from_json(file) -> str:
    try:
        json_content = json.load(file)
        text = json.dumps(json_content, indent=2)
        return text
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""


# Extract text from an XML file.
def extract_text_from_xml(file) -> str:
    try:
        xml_content = file.getvalue().decode("utf-8")
        root = ET.fromstring(xml_content)
        text = "\n".join([elem.text for elem in root.iter() if elem.text])
        return text
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""    


# Embed the texts and adds them to the FAISS index.
def add_texts_to_faiss(texts, embeddings):
    if not texts:
        return
    embeddings_list = embeddings.embed_documents(texts)
    embeddings_np = np.array(embeddings_list).astype('float32')
    st.session_state.faiss_index.add(embeddings_np)
    st.session_state.texts.extend(texts)
    st.session_state.metadatas.extend([{"source": "uploaded_file"}] * len(texts))


# Generate an answer based on the query and retrieved texts.
def generate_answer(client, query, verbose, embeddings, faiss_index, texts):
    # Embed the query
    query_embedding = embeddings.embed_query(query)
    query_embedding_np = np.array([query_embedding]).astype('float32')

    # Number of top documents to retrieve
    k = 20

    # Search FAISS index
    # TODO Why are duplicates included? Same phrase multiple times?
    distances, indices = faiss_index.search(query_embedding_np, k)

    if verbose:
        st.header("🔎 Search")
        distances
        indices

    # Build array of text snippets with lowest distances
    # TODO Combine snippets if they are adjacent (remove overlap)
    retrieved_texts = [texts[idx] for idx in indices[0] if idx < len(texts)]

    if not retrieved_texts:
        return "No relevant information found in the uploaded documents."

    # Combine retrieved texts
    context = "\n\n".join(retrieved_texts)

    # Prompt
    prompt = (
        "You are an AI assistant. Use the information provided below to answer the question. If you are uncertain, say so instead of guessing a response. If the question is given in a non-English language use that language in the response.\n\n"
        f"Context:\n{context}\n\n"
        f"Question: {query}\n"
    )

    if verbose:
        st.header("👨‍🎨 Prompt")
        st.text(prompt)

    # Generate the answer using OpenAI's GPT model
    return query_openai(client, prompt)


# Query OpenAI API
def query_openai(client, prompt):
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=4000,
            temperature=0.1,
            top_p=1,
            frequency_penalty=0,
            presence_penalty=0,
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error generating answer: {e}"


main()
