# https://numpy.org/
import numpy as np

# https://pypi.org/project/pandas/
import pandas as pd

# https://docs.python.org/3/library/os.html
import os

# https://docs.python.org/3/library/time.html
import time

# https://docs.python.org/3/library/datetime.html
from datetime import datetime

# https://docs.python.org/3/library/time.html
import html

# https://streamlit.io/
# https://docs.streamlit.io/develop/api-reference/text
import streamlit as st

# https://www.langchain.com/
from langchain.text_splitter import RecursiveCharacterTextSplitter

# https://python.langchain.com/docs/integrations/text_embedding/openai/
from langchain_openai import OpenAIEmbeddings

# https://faiss.ai/
# https://faiss.ai/cpp_api/struct/structfaiss_1_1IndexFlatL2.html
import faiss

# https://github.com/openai/openai-python
from openai import OpenAI

# https://pypdf2.readthedocs.io/en/3.x/
from PyPDF2 import PdfReader

# https://python-docx.readthedocs.io/en/latest/
import docx

# https://python-pptx.readthedocs.io/en/latest/
import pptx

# https://openpyxl.readthedocs.io/en/stable/
import openpyxl

# https://www.crummy.com/software/BeautifulSoup/bs4/doc/
from bs4 import BeautifulSoup

# https://docs.python.org/3/library/json.html
import json

# https://docs.python.org/3/library/xml.etree.elementtree.html
import xml.etree.ElementTree as ET

# Configuration
# https://docs.streamlit.io/develop/api-reference/configuration/config.toml

# Consider
#
# PDF to Markdown
# https://pypi.org/project/marker-pdf/
#
# FAISS on GPU
# https://github.com/facebookresearch/faiss/wiki/Running-on-GPUs

def main():
    # Application header
    st.set_page_config(page_title="Query Documents", page_icon="🙋‍♂️", layout="centered", menu_items={'About': 'https://abiro.com'})   
    st.title("📄 Query Documents")
    st.markdown("Upload your documents (PDF, Word, PowerPoint, Excel, Text, Markdown, HTML, XML or JSON), and ask questions based on their content. The system will only use information from the uploaded files to answer your queries. The provided API Key, set via OPENAI_API_KEY or manually, is only used for queries and not otherwise stored.")

    # AI configuration

    # OpenAI API key: via the OPENAI_API_KEY environment variable or input
    openai_api_key = os.getenv("OPENAI_API_KEY", "")
    if openai_api_key == "":
        st.header("🔑 OpenAI API Key")
        openai_api_key = st.text_input("Enter your API key (not shown):", type="password")
        st.button("Set")

    if openai_api_key == "":
        st.info("Enter a valid OpenAI API Key.", icon="ℹ️")
        return

    # SIDEBAR

    st.sidebar.header("Configuration")

    # Embedding model options with display names and values
    embedding_options = {
        "Embedding 3 small": "text-embedding-3-small",
        "Embedding 3 large": "text-embedding-3-large",
    }
    selected_embedding_display_default = 0
    selected_embedding_display = st.sidebar.selectbox("Embedding Model", list(embedding_options.keys()), index=selected_embedding_display_default)
    selected_embedding_value = embedding_options[selected_embedding_display]

    # RAG chunk size
    rag_chunk_size_default = 1000
    rag_chunk_size = st.sidebar.number_input("Embedding Chunk Size", min_value=100, max_value=10000, step=100, value=rag_chunk_size_default)

    # RAG chunk overlap
    rag_chunk_overlap_default = 100
    rag_chunk_overlap = st.sidebar.number_input("Embedding Chunk Overlap", min_value=10, max_value=1000, step=10, value=rag_chunk_overlap_default)    

    # GPT Model options with display names and values
    gpt_options = {
        "O1": "o1",
        "O1 Preview": "o1-preview",
        "O1 Mini": "o1-mini",
        "GPT-4.1": "gpt-4.1",
        "GPT-4.1 Mini": "gpt-4.1-mini",
        "GPT-4.1 Nano": "gpt-4.1-nano",
        "GPT-4o": "gpt-4o",
        "GPT-4o Mini": "gpt-4o-mini"
    }
    selected_gpt_display_default = 4
    selected_gpt_display = st.sidebar.selectbox("GPT Model", list(gpt_options.keys()), index=selected_gpt_display_default)
    selected_gpt_value = gpt_options[selected_gpt_display]

    # Set the temperature
    temperature_default = 0.1
    temperature = st.sidebar.slider("GPT Temperature", 0.0, 1.0, temperature_default, 0.1)

    # Input for instructions
    instructions_default = "You are an AI assistant. Use the information provided below to answer the question. If you are uncertain, say so instead of guessing a response."
    #  (If the question is given in a non-English language use that language in the response.)
    instructions = st.sidebar.text_area("GPT Instructions", instructions_default)

    st.session_state.verbose = st.sidebar.checkbox("Verbose", help="Check this if you want intermediate information and a timetrace")

    # Initialize session state for embeddings and FAISS index
    if 'embeddings' not in st.session_state:
        st.session_state.embeddings = OpenAIEmbeddings(model=selected_embedding_value, openai_api_key=openai_api_key)

    if 'faiss_index' not in st.session_state:
        init_faiss_index()

    if 'texts' not in st.session_state:
        st.session_state.texts = []

    if 'metadatas' not in st.session_state:
        st.session_state.metadatas = []

    if 'processed_files' not in st.session_state:
        st.session_state.processed_files = set()

    st.sidebar.header("Upload Files")

    uploaded_files = st.sidebar.file_uploader("Choose files", type=["pdf", "txt", "md", "docx", "pptx", "xlsx", "html", "htm", "json", "xml"], accept_multiple_files=True)
    if (len(uploaded_files) < len(st.session_state.processed_files)):
        st.session_state.processed_files = set()        
        init_faiss_index()

    timetrace_sample("Initialization done")

    if uploaded_files:
        for uploaded_file in uploaded_files:
            if uploaded_file.name not in st.session_state.processed_files:
                st.sidebar.write(f"📄 {uploaded_file.name}")
                # Extract text based on file type
                text = extract_text_from_file(uploaded_file)
                timetrace_sample("Text extraction done")

                if text:
                    if st.session_state.verbose:
                        st.header("📄 Parsed Text")
                        st.markdown(text)

                    # Split text into chunks
                    text_splitter = RecursiveCharacterTextSplitter(chunk_size=rag_chunk_size, chunk_overlap=rag_chunk_overlap)
                    chunks = text_splitter.split_text(text)
                    timetrace_sample("Text splitting done")

                    add_texts_to_faiss(chunks, st.session_state.embeddings)
                    timetrace_sample("FAISS adding done")

                    st.sidebar.success(f"Processed {uploaded_file.name} successfully!")
                    st.session_state.processed_files.add(uploaded_file.name)
                else:
                    st.sidebar.error(f"Failed to extract text from {uploaded_file.name}.")
            else:
                st.sidebar.write(f"📄 {uploaded_file.name} (already processed)")

    # Optional: Show number of indexed chunks
    st.sidebar.write(f"**Total Indexed Chunks:** {st.session_state.faiss_index.ntotal}")

    # Bail if no files uploaded
    if st.session_state.faiss_index.ntotal == 0:
        st.info("Please upload files to start querying.", icon="ℹ️")
        return

    # Query input
    st.header("🙋‍♂️ Question")
    user_query = st.text_input("Type your question here:")

    if st.button("Query") and user_query:
        with st.spinner("Querying..."):
            client = OpenAI(api_key=openai_api_key)
            answer = generate_answer(client, selected_gpt_value, temperature, instructions, user_query, st.session_state.verbose, st.session_state.embeddings, st.session_state.faiss_index, st.session_state.texts)
            timetrace_sample("GPT query done")
        
        st.header("👨‍🏫 Answer")
        st.write(answer)


def extract_text_from_file(file) -> str:
    try:
        match file.type:
            case "application/pdf":
                return extract_text_from_pdf(file)
            case "text/plain":
                return extract_text_from_plain_text(file)
            case "text/markdown":
                return extract_text_from_markdown(file)
            case "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                return extract_text_from_docx(file)
            case "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return extract_text_from_pptx(file)
            case "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                return extract_text_from_xlsx(file)
            case "text/html":
                return extract_text_from_html(file)
            case "application/json":
                return extract_text_from_json(file)
            case "application/xml" | "text/xml":
                return extract_text_from_xml(file)
            case "application/octet-stream":
                return extract_text_from_octet_stream(file)
            case _:
                return ""
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""

def extract_text_from_plain_text(file) -> str:
    return f"```\n{file.getvalue().decode('utf-8')}\n```"

def extract_text_from_markdown(file) -> str:
    return file.getvalue().decode("utf-8")

def extract_text_from_pdf(file) -> str:
    try:
        pdf = PdfReader(file)
        text = ""
        for page_num, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text:
                text += page_text
        return f"```\n{text}\n```"
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""

def extract_text_from_docx(file) -> str:
    try:
        doc = docx.Document(file)
        text = "\n".join([f"### {para.text}" if para.style.name.startswith('Heading') else para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""

def extract_text_from_pptx(file) -> str:
    try:
        presentation = pptx.Presentation(file)
        text = ""
        for slide_num, slide in enumerate(presentation.slides, start=1):
            text += f"## Slide {slide_num}\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += f"{shape.text}\n"
        return text
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""

def extract_text_from_xlsx(file) -> str:
    try:
        workbook = openpyxl.load_workbook(file, data_only=True)
        text = ""
        for sheet in workbook.sheetnames:
            text += f"## {sheet}\n"
            worksheet = workbook[sheet]
            for row in worksheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                text += "| " + " | ".join(row_text) + " |\n"
        return text
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""

def extract_text_from_html(file) -> str:
    try:
        html_content = file.getvalue().decode("utf-8")
        soup = BeautifulSoup(html_content, "html.parser")
        text = soup.get_text(separator="\n")
        return f"```\n{text}\n```"
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""

def extract_text_from_json(file) -> str:
    try:
        json_content = json.load(file)
        text = json.dumps(json_content, indent=2)
        return f"```json\n{text}\n```"
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""

def extract_text_from_xml(file) -> str:
    try:
        xml_content = file.getvalue().decode("utf-8")
        root = ET.fromstring(xml_content)
        text = "\n".join([elem.text for elem in root.iter() if elem.text])
        return f"```\n{text}\n```"
    except Exception as e:
        st.error(f"Error reading {file.name}: {e}")
        return ""

def extract_text_from_octet_stream(file) -> str:
    return f"```\n{file.getvalue().decode('utf-8')}\n```"


# Initialize FAISS index with the dimensionality from embeddings
def init_faiss_index():
    embedding_test = st.session_state.embeddings.embed_query("test")
    embedding_dim = len(embedding_test)
    st.session_state.faiss_index = faiss.IndexFlatL2(embedding_dim)


# Embed the texts and add them to the FAISS index.
def add_texts_to_faiss(texts, embeddings):
    if not texts:
        return
    embeddings_list = embeddings.embed_documents(texts)
    embeddings_np = np.array(embeddings_list).astype('float32')
    st.session_state.faiss_index.add(embeddings_np)
    st.session_state.texts.extend(texts)
    st.session_state.metadatas.extend([{"source": "uploaded_file"}] * len(texts))


# Generate an answer based on the query and retrieved texts.
def generate_answer(client, model, temperature, instructions, query, verbose, embeddings, faiss_index, texts):
    # Embed the query
    query_embedding = embeddings.embed_query(query)
    query_embedding_np = np.array([query_embedding]).astype('float32')
    timetrace_sample("Embedding query done")    

    # Number of top documents to retrieve
    k = 20

    # Search FAISS index
    # TODO Why are duplicates included? Same phrase multiple times?
    distances, indices = faiss_index.search(query_embedding_np, k)
    timetrace_sample("FAISS search done")

    if verbose:
        st.header("🔎 Search")
        distances
        indices

    # Build array of text snippets with lowest distances
    # TODO Combine snippets if they are adjacent (remove overlap)
    retrieved_texts = [texts[idx] for idx in indices[0] if idx < len(texts)]
    timetrace_sample("Text collation done")
    if not retrieved_texts:
        return "No relevant information found in the uploaded documents."

    # Combine retrieved texts
    context = "\n\n".join(retrieved_texts)

    # Prompt
    prompt = (
        instructions + "\n\n"
        f"Context:\n{context}\n\n"
        f"Question: {query}\n"
    )

    if verbose:
        st.header("👨‍🎨 Prompt")
        st.text(prompt)

    # Generate the answer using OpenAI's GPT model
    return query_openai(client, model, temperature, prompt)


# Query OpenAI API
def query_openai(client, model, temperature, prompt):
    try:
        response = client.chat.completions.create(
            model=model,
            temperature=temperature,
            top_p=1,
            frequency_penalty=0,
            presence_penalty=0,
            messages=[{"role": "user", "content": prompt}]
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error generating answer: {e}"


# Timetrace
timetrace_previous = 0
timetrace_samples = []

def timetrace_init():
    global timetrace_previous
    global timetrace_samples

    timetrace_previous = time.time()
    timetrace_samples = []

def timetrace_sample(note):
    global timetrace_previous
    global timetrace_samples

    current = time.time()
    timetrace_samples.append(str(round(current - timetrace_previous, 3)) + ": " + note)
    timetrace_previous = current


def timetrace_show():
    global timetrace_samples

    st.header("⏰ Timetrace")

    st.write("Time in seconds.")

    st.markdown('\n'.join([f"* {item}" for item in timetrace_samples]))

    
# Launch

timetrace_init()
st.session_state.verbose = False

timetrace_sample("Processing starting")

main()

timetrace_sample("Processing done")

if st.session_state.verbose:
    timetrace_show()

# Footer

current_year = datetime.now().year

st.html(f"""
<style>
    .footer {{

    }}
    .footer-link {{
        color: inherit;
    }}
</style>            

<hr/>
            
<footer class="footer">
    © {current_year} <a href="https://abiro.com" target="_blank" class="footer-link">Abiro</a> All rights reserved. 
    <a href="https://abiro.com/about/privacy-policy/" target="_blank" class="footer-link">Privacy Policy</a>. 
    <a href="https://apps.abiro.com/" target="_blank" class="footer-link">Other Applications</a>.
</footer>
""")
